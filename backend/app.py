"""
Backend API for fetching LangSmith traces and rendering PPTX natively with Konva.

Key fixes vs previous version:
- Correctly reads font styling from paragraph.font when runs are missing
- Extracts text frame margins (padding) + vertical anchor (vAlign) for better Konva layout
- Uses robust alignment parsing (PP_ALIGN constants)
- Improves background/fill extraction and defaults
"""

from dotenv import load_dotenv
load_dotenv()

import os
import ast
import base64
from io import BytesIO
from typing import Optional, Any

from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from langsmith import Client

# PPTX parsing and building
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


app = FastAPI(title="Slide Viewer API")

# CORS for React frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173", "http://localhost:3000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize LangSmith client
ls_client = Client()


# ============================================================================
# PPTX PARSING - Extract slide data as JSON for Konva rendering
# ============================================================================

PIXELS_PER_INCH = 96
EMU_PER_INCH = 914400


def emu_to_pixels(emu: int) -> float:
    """Convert EMUs to pixels (96 DPI)."""
    inches = emu / EMU_PER_INCH
    return inches * PIXELS_PER_INCH


def pixels_to_emu(pixels: float) -> int:
    """Convert pixels to EMUs (96 DPI)."""
    inches = pixels / PIXELS_PER_INCH
    return int(inches * EMU_PER_INCH)


def rgb_to_hex(obj: Any) -> Optional[str]:
    """
    Convert pptx color objects to hex string.
    Handles RGBColor and ColorFormat-ish objects.
    """
    if obj is None:
        return None

    try:
        # If this is a ColorFormat, it may have .rgb which is RGBColor or None
        if hasattr(obj, "rgb") and obj.rgb is not None:
            rgb = obj.rgb
            if isinstance(rgb, RGBColor):
                return f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"
            # Sometimes .rgb can be bytes
            if isinstance(rgb, (bytes, bytearray)) and len(rgb) == 3:
                return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
        # If it's already an RGBColor
        if isinstance(obj, RGBColor):
            return f"#{obj.red:02x}{obj.green:02x}{obj.blue:02x}"
    except Exception:
        return None

    return None


def hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    """Convert hex string to (r,g,b)."""
    hex_color = (hex_color or "").lstrip("#")
    if len(hex_color) != 6:
        return (255, 255, 255)
    return tuple(int(hex_color[i:i + 2], 16) for i in (0, 2, 4))


def get_background_color(slide) -> str:
    """Extract background color from slide."""
    try:
        fill = slide.background.fill
        # If the fill is solid, read fore_color
        if fill and fill.type is not None:
            color = rgb_to_hex(fill.fore_color)
            if color:
                return color
    except Exception:
        pass
    return "#0f172a"


def get_text_alignment(alignment) -> str:
    """Convert PP_ALIGN to string."""
    if alignment is None:
        return "left"
    if alignment == PP_ALIGN.CENTER:
        return "center"
    if alignment == PP_ALIGN.RIGHT:
        return "right"
    if alignment == PP_ALIGN.JUSTIFY:
        return "justify"
    return "left"


def get_vertical_alignment(anchor) -> str:
    """
    Convert MSO_ANCHOR to string for Konva.
    PowerPoint anchor is vertical alignment inside textbox.
    """
    if anchor is None:
        return "top"
    if anchor == MSO_ANCHOR.MIDDLE:
        return "middle"
    if anchor == MSO_ANCHOR.BOTTOM:
        return "bottom"
    return "top"


def _font_to_style(font) -> dict:
    """Extract font properties from a pptx font object (run.font or paragraph.font)."""
    out: dict = {}

    # size
    try:
        if getattr(font, "size", None) is not None:
            out["fontSize"] = float(font.size.pt)
    except Exception:
        pass

    # name/family
    try:
        name = getattr(font, "name", None)
        if name:
            out["fontFamily"] = name
    except Exception:
        pass

    # color
    try:
        color = rgb_to_hex(getattr(font, "color", None))
        if color:
            out["fill"] = color
    except Exception:
        pass

    # weight/style
    bold = bool(getattr(font, "bold", False))
    italic = bool(getattr(font, "italic", False))
    if bold or italic:
        parts = []
        if bold:
            parts.append("bold")
        if italic:
            parts.append("italic")
        out["fontStyle"] = " ".join(parts)
    else:
        out["fontStyle"] = "normal"

    return out


def _safe_float(v, default: float = 0.0) -> float:
    try:
        return float(v)
    except Exception:
        return default


def parse_shape(shape) -> Optional[dict]:
    """
    Parse a single shape into JSON for Konva.
    Adds extra fields for better text layout: padding, vAlign, lineHeight.
    """
    try:
        shape_data: dict = {
            "id": str(shape.shape_id),
            "name": getattr(shape, "name", ""),
            "x": emu_to_pixels(shape.left),
            "y": emu_to_pixels(shape.top),
            "width": emu_to_pixels(shape.width),
            "height": emu_to_pixels(shape.height),
            "type": "rect",
            "draggable": True,
        }

        # TEXTBOX / TEXT FRAME
        if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
            tf = shape.text_frame
            text_content = (tf.text or "").strip()
            if text_content:
                shape_data["type"] = "text"
                shape_data["text"] = tf.text or ""

                # Textbox internal padding (important for Konva)
                try:
                    shape_data["padding"] = {
                        "left": emu_to_pixels(tf.margin_left),
                        "right": emu_to_pixels(tf.margin_right),
                        "top": emu_to_pixels(tf.margin_top),
                        "bottom": emu_to_pixels(tf.margin_bottom),
                    }
                except Exception:
                    shape_data["padding"] = {"left": 0, "right": 0, "top": 0, "bottom": 0}

                # Vertical anchor (top/middle/bottom)
                try:
                    shape_data["vAlign"] = get_vertical_alignment(tf.vertical_anchor)
                except Exception:
                    shape_data["vAlign"] = "top"

                # Word wrap
                try:
                    shape_data["wrap"] = bool(tf.word_wrap)
                except Exception:
                    shape_data["wrap"] = True

                # Paragraph alignment + style
                if tf.paragraphs:
                    para = tf.paragraphs[0]
                    shape_data["align"] = get_text_alignment(para.alignment)

                    # Prefer run font if present, else paragraph font
                    if para.runs:
                        style = _font_to_style(para.runs[0].font)
                    else:
                        style = _font_to_style(para.font)

                    # Defaults
                    shape_data["fontSize"] = style.get("fontSize", 16.0)
                    shape_data["fill"] = style.get("fill", "#ffffff")
                    shape_data["fontStyle"] = style.get("fontStyle", "normal")
                    if "fontFamily" in style:
                        shape_data["fontFamily"] = style["fontFamily"]

                    # Line spacing (Konva uses lineHeight multiplier)
                    # pptx line_spacing is typically points or None; space_after/before too.
                    # We'll approximate a reasonable lineHeight.
                    try:
                        # python-pptx exposes para.line_spacing as Length or float-ish sometimes
                        ls = para.line_spacing
                        if ls is None:
                            shape_data["lineHeight"] = 1.1
                        else:
                            # If it's an absolute point value, approximate multiplier by / fontSize
                            if hasattr(ls, "pt"):
                                fs = shape_data["fontSize"] or 16.0
                                shape_data["lineHeight"] = max(1.0, _safe_float(ls.pt, fs) / fs)
                            else:
                                # Might already be a multiplier
                                shape_data["lineHeight"] = max(1.0, _safe_float(ls, 1.1))
                    except Exception:
                        shape_data["lineHeight"] = 1.1
                else:
                    shape_data["align"] = "left"
                    shape_data["fontSize"] = 16.0
                    shape_data["fill"] = "#ffffff"
                    shape_data["fontStyle"] = "normal"
                    shape_data["lineHeight"] = 1.1

                # Text should not have a rect fill by default
                return shape_data

        # AUTO SHAPES (RECTANGLES, ETC.)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            shape_data["type"] = "rect"
            # Fill color
            try:
                fill = shape.fill
                if fill and fill.type is not None:
                    color = rgb_to_hex(fill.fore_color)
                    if color:
                        shape_data["fill"] = color
            except Exception:
                pass

            if "fill" not in shape_data:
                # default card fill
                shape_data["fill"] = "#1e293b"

            # Borders: if you want, you can parse line color/width here
            return shape_data

        # Ignore other shape types (pictures, charts) for now
        return None

    except Exception as e:
        print(f"Error parsing shape: {e}")
        return None


def parse_pptx_bytes(pptx_bytes: bytes) -> dict:
    """Parse PPTX bytes into JSON structure for Konva rendering."""
    try:
        prs = Presentation(BytesIO(pptx_bytes))

        result = {
            "width": emu_to_pixels(prs.slide_width),
            "height": emu_to_pixels(prs.slide_height),
            "slides": []
        }

        for idx, slide in enumerate(prs.slides):
            slide_data = {
                "id": f"slide-{idx+1}",
                "index": idx,
                "backgroundColor": get_background_color(slide),
                "shapes": []
            }

            for shape in slide.shapes:
                parsed = parse_shape(shape)
                if parsed:
                    slide_data["shapes"].append(parsed)

            result["slides"].append(slide_data)

        return result

    except Exception as e:
        print(f"Error parsing PPTX: {e}")
        import traceback
        traceback.print_exc()
        return {"error": str(e), "slides": [], "width": 1280, "height": 720}


# ============================================================================
# PPTX BUILDING - Rebuild PPTX from edited JSON
# ============================================================================

def build_pptx_from_json(presentation_data: dict) -> bytes:
    """
    Build a PPTX file from JSON presentation data.
    Supports background + rect + text with padding/alignment.
    """
    prs = Presentation()

    width_px = float(presentation_data.get("width", 1280))
    height_px = float(presentation_data.get("height", 720))
    prs.slide_width = pixels_to_emu(width_px)
    prs.slide_height = pixels_to_emu(height_px)

    blank_layout = prs.slide_layouts[6]

    for slide_data in presentation_data.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)

        # Background
        bg_hex = slide_data.get("backgroundColor", "#0f172a")
        r, g, b = hex_to_rgb(bg_hex)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)

        for shape_data in slide_data.get("shapes", []):
            shape_type = shape_data.get("type", "rect")
            x = pixels_to_emu(shape_data.get("x", 0))
            y = pixels_to_emu(shape_data.get("y", 0))
            w = pixels_to_emu(shape_data.get("width", 100))
            h = pixels_to_emu(shape_data.get("height", 50))

            if shape_type == "rect":
                rect = slide.shapes.add_shape(1, x, y, w, h)
                fill_hex = shape_data.get("fill", "#1e293b")
                rr, gg, bb = hex_to_rgb(fill_hex)
                rect.fill.solid()
                rect.fill.fore_color.rgb = RGBColor(rr, gg, bb)
                rect.line.fill.background()

            elif shape_type == "text":
                textbox = slide.shapes.add_textbox(x, y, w, h)
                tf = textbox.text_frame
                tf.word_wrap = bool(shape_data.get("wrap", True))

                # padding (text frame margins)
                pad = shape_data.get("padding") or {}
                try:
                    tf.margin_left = pixels_to_emu(pad.get("left", 0))
                    tf.margin_right = pixels_to_emu(pad.get("right", 0))
                    tf.margin_top = pixels_to_emu(pad.get("top", 0))
                    tf.margin_bottom = pixels_to_emu(pad.get("bottom", 0))
                except Exception:
                    pass

                # vertical anchor
                v_align = shape_data.get("vAlign", "top")
                try:
                    if v_align == "middle":
                        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    elif v_align == "bottom":
                        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
                    else:
                        tf.vertical_anchor = MSO_ANCHOR.TOP
                except Exception:
                    pass

                p = tf.paragraphs[0]
                p.text = shape_data.get("text", "")

                # style
                font_size = float(shape_data.get("fontSize", 16))
                p.font.size = Pt(font_size)

                font_family = shape_data.get("fontFamily")
                if font_family:
                    p.font.name = font_family

                fill_hex = shape_data.get("fill", "#ffffff")
                rr, gg, bb = hex_to_rgb(fill_hex)
                p.font.color.rgb = RGBColor(rr, gg, bb)

                font_style = shape_data.get("fontStyle", "normal")
                p.font.bold = "bold" in font_style
                p.font.italic = "italic" in font_style

                align = shape_data.get("align", "left")
                if align == "center":
                    p.alignment = PP_ALIGN.CENTER
                elif align == "right":
                    p.alignment = PP_ALIGN.RIGHT
                elif align == "justify":
                    p.alignment = PP_ALIGN.JUSTIFY
                else:
                    p.alignment = PP_ALIGN.LEFT

                # line spacing (approx)
                try:
                    lh = float(shape_data.get("lineHeight", 1.1))
                    # Convert multiplier to approximate points
                    # e.g., 1.2 * fontSize points
                    p.line_spacing = Pt(max(1.0, lh) * font_size)
                except Exception:
                    pass

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


# ============================================================================
# PYDANTIC MODELS
# ============================================================================

class PaddingData(BaseModel):
    left: float = 0
    right: float = 0
    top: float = 0
    bottom: float = 0


class ShapeData(BaseModel):
    id: str
    name: Optional[str] = None
    x: float
    y: float
    width: float
    height: float
    type: str
    text: Optional[str] = None
    fontSize: Optional[float] = None
    fontFamily: Optional[str] = None
    fill: Optional[str] = None
    fontStyle: Optional[str] = None
    align: Optional[str] = None
    vAlign: Optional[str] = None
    lineHeight: Optional[float] = None
    wrap: Optional[bool] = True
    padding: Optional[PaddingData] = None
    draggable: Optional[bool] = True


class SlideData(BaseModel):
    id: str
    index: int
    backgroundColor: str
    shapes: list[ShapeData]


class PresentationData(BaseModel):
    width: float
    height: float
    slides: list[SlideData]
    error: Optional[str] = None


class TraceSlide(BaseModel):
    trace_id: str
    trace_name: str
    created_at: str
    pptx_base64: Optional[str] = None
    presentation: Optional[PresentationData] = None
    error: Optional[str] = None


class TracesResponse(BaseModel):
    traces: list[TraceSlide]


class SavePresentationRequest(BaseModel):
    presentation: PresentationData


class SavePresentationResponse(BaseModel):
    pptx_base64: str
    success: bool


# ============================================================================
# API ENDPOINTS
# ============================================================================

def extract_pptx_from_trace(trace_id: str) -> Optional[bytes]:
    """Extract PPTX bytes from a trace's finalize_presentation tool call."""
    project_name = os.getenv("LANGSMITH_PROJECT", "default")
    try:
        runs = list(ls_client.list_runs(
            project_name=project_name,
            trace_id=trace_id,
        ))

        for run in runs:
            if run.name == "finalize_presentation" and run.outputs:
                output = run.outputs.get("output")
                if not output:
                    continue

                content = output.get("content")
                if not content:
                    continue

                # content is often like "b'...'"
                if isinstance(content, bytes):
                    return content

                if isinstance(content, str):
                    try:
                        pptx_bytes = ast.literal_eval(content)
                        if isinstance(pptx_bytes, (bytes, bytearray)):
                            print(f"Extracted {len(pptx_bytes)} bytes from trace")
                            return bytes(pptx_bytes)
                    except Exception:
                        # Sometimes it might already be base64 or not literal-evaluable
                        pass

        print(f"No finalize_presentation output found in trace {trace_id}")
        return None

    except Exception as e:
        print(f"Error extracting PPTX from trace {trace_id}: {e}")
        import traceback
        traceback.print_exc()
        return None


@app.get("/api/traces", response_model=TracesResponse)
async def get_recent_traces():
    """Get the last 3 traces with their PPTX outputs, parsed for Konva rendering."""
    project_name = os.getenv("LANGSMITH_PROJECT", "default")

    try:
        root_runs = list(ls_client.list_runs(
            project_name=project_name,
            is_root=True,
            limit=3,
        ))

        result_traces: list[TraceSlide] = []
        for run in root_runs:
            trace_slide = TraceSlide(
                trace_id=str(run.trace_id),
                trace_name=run.name or "Unnamed",
                created_at=run.start_time.isoformat() if run.start_time else "",
            )

            pptx_bytes = extract_pptx_from_trace(str(run.trace_id))
            if pptx_bytes:
                parsed = parse_pptx_bytes(pptx_bytes)

                if parsed.get("slides"):
                    trace_slide.presentation = PresentationData(**parsed)
                elif parsed.get("error"):
                    trace_slide.error = parsed["error"]

                trace_slide.pptx_base64 = base64.b64encode(pptx_bytes).decode()
            else:
                trace_slide.error = "No PPTX output found in trace"

            result_traces.append(trace_slide)

        return TracesResponse(traces=result_traces)

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/save-presentation", response_model=SavePresentationResponse)
async def save_presentation(request: SavePresentationRequest):
    """Rebuild PPTX from edited presentation data and return as base64."""
    try:
        pptx_bytes = build_pptx_from_json(request.presentation.model_dump())
        pptx_base64 = base64.b64encode(pptx_bytes).decode()
        return SavePresentationResponse(pptx_base64=pptx_base64, success=True)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/health")
async def health():
    return {
        "status": "ok",
        "langsmith_project": os.getenv("LANGSMITH_PROJECT", "default")
    }


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
